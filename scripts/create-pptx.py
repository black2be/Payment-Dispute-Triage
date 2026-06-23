#!/usr/bin/env python3
"""Generate a PowerPoint presentation for the Payment Dispute Triage demo."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Standard Bank colours
SB_BLUE = RGBColor(0x00, 0x33, 0xA1)
SB_GOLD = RGBColor(0xC8, 0xA4, 0x15)
SB_DARK = RGBColor(0x00, 0x23, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x4B, 0x55, 0x63)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)


def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Blue background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = SB_BLUE

    # Logo
    slide.shapes.add_picture("docs/sbsa logo.jpeg", Inches(5.4), Inches(0.8), height=Inches(1.5))

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(1))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(20)
    p2.font.color.rgb = SB_GOLD
    p2.alignment = PP_ALIGN.CENTER


def add_section_slide(title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = SB_DARK

    # Logo top-right
    slide.shapes.add_picture("docs/sbsa logo.jpeg", Inches(11.3), Inches(0.4), height=Inches(0.9))

    txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


def add_content_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title bar
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = SB_BLUE

    # Content
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = GRAY
        p.space_after = Pt(8)


def add_table_slide(title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = SB_BLUE

    # Table
    cols = len(headers)
    tbl_rows = len(rows) + 1
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(12.3)
    height = Inches(0.4) * tbl_rows

    table_shape = slide.shapes.add_table(tbl_rows, cols, left, top, width, height)
    table = table_shape.table

    # Headers
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = SB_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE

    # Rows
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = GRAY


# ─── SLIDES ───

# Slide 1: Title
add_title_slide(
    "Payment Dispute Triage",
    "AI SDLC — Digital Platforms\nEngineering Community of Practice"
)

# Slide 2: Agenda
add_content_slide("Agenda", [
    "1. Specifications — How roles contributed",
    "2. Harness — Steering files, hooks, and specs",
    "3. Demo — Live prototype demonstration",
    "4. What We Learnt — Rethinking software development",
])

# ─── SECTION 1: SPECIFICATIONS ───
add_section_slide("1. Specifications")

add_content_slide("The Challenge", [
    "Build a rules-based payment dispute triage prototype",
    "Help operations users determine the next best action for a dispute",
    "Transparent, reproducible, and explainable decisions",
    "Mock-only, no external integrations — focused on the decision logic",
])

add_table_slide("Role Contributions", 
    ["Role", "Contribution", "Artefact"],
    [
        ["Product Owner", "Vision, scope, 4 actions, guardrails", "product.md, requirements.md"],
        ["Architect", "3-tier design, layering, data model", "design.md §1–2, architecture.md"],
        ["API Designer", "REST endpoints, validation, error shapes", "design.md §4, api-spec.md"],
        ["UI/UX Designer", "Single-screen flow, palette, accessibility", "design.md §5, ui-spec.md"],
        ["Rules Engineer", "6 ordered rules, priority logic, 7 worked examples", "design.md §3"],
        ["Test Engineer", "51 test cases, property tests P1–P7", "test-plan.md, testing-standards.md"],
        ["Harness Engineer", "Governance hooks, CI/CD, compliance", "governance.md, .kiro/hooks/"],
    ]
)

add_content_slide("Key Design Decisions", [
    "• Pure TypeScript rules engine — deterministic, no I/O",
    "• First-match-wins precedence (R1→R6)",
    "• Inject 'today' — no Date.now() inside the engine",
    "• Code-form enums in API/DB, label-form in UI",
    "• Persist recommendation even if display unavailable (REQ-05.5)",
    "• 7 worked examples as deterministic oracle for testing",
])

# ─── SECTION 2: HARNESS ───
add_section_slide("2. The Harness")

add_table_slide("Steering Files (.kiro/steering/)",
    ["File", "Purpose"],
    [
        ["00-instructions.md", "Prime directive + guardrails (read first)"],
        ["product.md", "What/why/scope — the user journey"],
        ["tech.md", "Stack: React 18, Vite, Tailwind, Vitest, fast-check"],
        ["structure.md", "Directory layout, naming, layering rule"],
        ["conventions.md", "Coding patterns, rule array, inject today"],
        ["governance.md", "AI SDLC — mock-only, approved MCP, no PII"],
        ["testing-standards.md", "P1–P7, worked examples, boundaries"],
        ["api-standards.md", "REST conventions, validation, error shape"],
    ]
)

add_content_slide("How Steering Guided Kiro", [
    "• testing-standards.md auto-loaded when test files are open",
    "• Governance rules prevented PII/secrets in generated code",
    "• Structure rules enforced: engine never imports React",
    "• Conventions ensured thresholds in constants.ts — no magic numbers",
    "• Kiro's first attempt was usually correct because it had team decisions",
])

add_table_slide("Hooks (.kiro/hooks/)",
    ["Hook", "Trigger", "Action"],
    [
        ["lint-on-save", "File edited (*.ts, *.tsx)", "Run ESLint"],
        ["typecheck-on-save", "File edited (*.ts, *.tsx)", "Run tsc --noEmit"],
        ["test-on-change", "File edited (*.test.ts)", "Run Vitest"],
    ]
)

add_content_slide("Specs (.kiro/specs/payment-triage/)", [
    "• requirements.md — REQ-01 to REQ-07 with acceptance criteria",
    "• design.md — Architecture, data model, rules, API, UI",
    "• tasks.md — Ordered implementation plan (engine first, UI second)",
    "",
    "The spec drove task execution top-to-bottom.",
    "Each task names the REQ it satisfies — full traceability.",
])

# ─── SECTION 3: DEMO ───
add_section_slide("3. Live Demo")

add_content_slide("Demo Scenarios", [
    "1. Happy path — demonstrate all 4 recommended actions:",
    "   Resolve Immediately | Investigate Further | Escalate | Refer",
    "",
    "2. Validation — empty form, future date rejection",
    "",
    "3. Reference lookup — auto-populate vs manual entry",
    "",
    "4. Rule transparency — all 6 rules visible, triggered one highlighted",
    "",
    "5. Single screen — no navigation, summary + recommendation together",
])

add_content_slide("Demo URL & Test Evidence", [
    "Live app: http://52.19.75.103:8081",
    "",
    "Test results:",
    "• 13 E2E tests with video recordings (Playwright)",
    "• 27 unit/property tests (Vitest + fast-check)",
    "• 9 API integration tests (supertest)",
    "• All 51 test cases from the test plan covered",
    "",
    "Run: npx playwright show-report  (HTML report with videos)",
])

# ─── SECTION 4: WHAT WE LEARNT ───
add_section_slide("4. What We Learnt")

add_table_slide("Mindset Shift",
    ["Before", "After"],
    [
        ["Write code, then document", "Document decisions, then code flows from them"],
        ["AI generates code, I review", "AI generates code within my team's constraints"],
        ["Tests prove code works", "Properties prove the specification is upheld"],
        ["Governance = manual checklist", "Governance = automated hooks + CI gates"],
        ["One person knows architecture", "Steering files mean everyone (and AI) knows"],
    ]
)

add_content_slide("Key Insights", [
    "• Writing the spec first didn't slow us down — it eliminated bugs before they existed",
    "• The harness transforms AI from 'smart autocomplete' into 'team member that follows standards'",
    "• Property-based testing found edge cases we never would have written manually",
    "• Automated guardrails aren't bureaucracy — they're confidence",
    "• The 7 worked examples caught 3 logic errors before any code existed",
])

add_content_slide("The Bottom Line", [
    "",
    "The AI SDLC harness transforms AI coding assistants",
    "from 'smart autocomplete'",
    "into 'team members that follow your standards.'",
    "",
    "Prototype: ~2 hours of active work.",
    "Without the harness: same time, more rework,",
    "inconsistency, and governance gaps.",
])

# Final slide
add_title_slide(
    "Thank You",
    "Questions?\n\nPayment Dispute Triage — Digital Platforms"
)

# Save
prs.save("docs/Payment-Dispute-Triage-Presentation.pptx")
print("Created: docs/Payment-Dispute-Triage-Presentation.pptx")
